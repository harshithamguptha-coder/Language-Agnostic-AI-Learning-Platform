import os
from dataclasses import dataclass
from pathlib import Path
from uuid import uuid4

import fitz
import pytesseract
from dotenv import load_dotenv
from PIL import Image, ImageEnhance, ImageFilter, ImageOps
from docx import Document
from pptx import Presentation
from fastapi import UploadFile

BACKEND_DIR = Path(__file__).resolve().parents[1]
load_dotenv(BACKEND_DIR / ".env")
load_dotenv()

TESSERACT_CMD = os.getenv("TESSERACT_CMD")
if TESSERACT_CMD:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD

OCR_LANGUAGES = os.getenv("OCR_LANGUAGES", "eng+hin+kan")
OCR_MIN_TEXT_CHARS = int(os.getenv("OCR_MIN_TEXT_CHARS", "40"))
PDF_OCR_DPI = int(os.getenv("PDF_OCR_DPI", "220"))
SUPPORTED_TEXT_TYPES = [".txt", ".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg", ".bmp", ".webp", ".tif", ".tiff"]


@dataclass
class ExtractionResult:
    text: str
    file_name: str
    file_type: str
    extraction_method: str
    page_count: int | None = None
    ocr_pages: int = 0
    average_confidence: float | None = None

    @property
    def text_length(self) -> int:
        return len(self.text)

    def metadata(self) -> dict:
        return {
            "file_name": self.file_name,
            "file_type": self.file_type,
            "extraction_method": self.extraction_method,
            "page_count": self.page_count,
            "ocr_pages": self.ocr_pages,
            "average_confidence": self.average_confidence,
            "text_length": self.text_length,
        }


def save_upload_file(upload_file: UploadFile, target_dir: str) -> str:
    extension = os.path.splitext(upload_file.filename)[1].lower()
    if extension not in SUPPORTED_TEXT_TYPES:
        raise ValueError("Unsupported file type")
    os.makedirs(target_dir, exist_ok=True)
    safe_name = f"{uuid4().hex}{extension}"
    destination = os.path.join(target_dir, safe_name)
    with open(destination, "wb") as buffer:
        buffer.write(upload_file.file.read())
    return destination


def extract_text_from_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as file:
        return file.read().strip()


def extract_text_from_pdf(path: str) -> str:
    content = []
    with fitz.open(path) as doc:
        for page in doc:
            content.append(page.get_text())
    return "\n".join(content).strip()


def _prepare_image_for_ocr(image: Image.Image) -> Image.Image:
    image = ImageOps.exif_transpose(image).convert("RGB")
    image = ImageOps.grayscale(image)
    image = ImageOps.autocontrast(image)
    image = ImageEnhance.Contrast(image).enhance(1.8)
    image = image.filter(ImageFilter.MedianFilter(size=3))

    width, height = image.size
    if max(width, height) < 1800:
        image = image.resize((width * 2, height * 2), Image.Resampling.LANCZOS)

    return image.filter(ImageFilter.SHARPEN)


def _read_confidence_values(ocr_data: dict) -> list[float]:
    confidence_values = []
    for raw_confidence in ocr_data.get("conf", []):
        try:
            confidence = float(raw_confidence)
        except (TypeError, ValueError):
            continue
        if confidence >= 0:
            confidence_values.append(confidence)
    return confidence_values


def _run_tesseract(image: Image.Image, config: str) -> tuple[str, list[float]]:
    prepared_image = _prepare_image_for_ocr(image)
    try:
        data = pytesseract.image_to_data(
            prepared_image,
            lang=OCR_LANGUAGES,
            config=config,
            output_type=pytesseract.Output.DICT,
        )
    except pytesseract.TesseractError:
        if OCR_LANGUAGES == "eng":
            raise
        data = pytesseract.image_to_data(
            prepared_image,
            lang="eng",
            config=config,
            output_type=pytesseract.Output.DICT,
        )

    words = [word.strip() for word in data.get("text", []) if word and word.strip()]
    return " ".join(words).strip(), _read_confidence_values(data)


def _ocr_image(image: Image.Image) -> tuple[str, float | None]:
    candidates = []
    for config in ("--oem 3 --psm 6", "--oem 3 --psm 11"):
        text, confidence_values = _run_tesseract(image, config)
        average_confidence = (
            round(sum(confidence_values) / len(confidence_values), 2)
            if confidence_values
            else None
        )
        candidates.append((text, average_confidence))

    return max(
        candidates,
        key=lambda candidate: (
            candidate[1] if candidate[1] is not None else -1,
            len(candidate[0]),
        ),
    )


def extract_pdf_with_ocr(path: str) -> tuple[str, int, int, float | None, str]:
    pages = []
    ocr_pages = 0
    confidence_values = []
    zoom = PDF_OCR_DPI / 72
    matrix = fitz.Matrix(zoom, zoom)

    with fitz.open(path) as doc:
        page_count = len(doc)
        for page_number, page in enumerate(doc, start=1):
            page_text = page.get_text().strip()
            if len(page_text) >= OCR_MIN_TEXT_CHARS:
                pages.append(f"[Page {page_number}]\n{page_text}")
                continue

            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
            ocr_text, confidence = _ocr_image(image)
            if ocr_text:
                pages.append(f"[Page {page_number} OCR]\n{ocr_text}")
                ocr_pages += 1
                if confidence is not None:
                    confidence_values.append(confidence)

    average_confidence = (
        round(sum(confidence_values) / len(confidence_values), 2)
        if confidence_values
        else None
    )
    method = "pdf_text_and_ocr" if ocr_pages else "pdf_text"
    return "\n\n".join(pages).strip(), page_count, ocr_pages, average_confidence, method


def extract_text_from_docx(path: str) -> str:
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs).strip()


def extract_text_from_pptx(path: str) -> str:
    presentation = Presentation(path)
    extracted = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    extracted.append(text)
    return "\n".join(extracted).strip()


def extract_text_from_image(path: str) -> str:
    image = Image.open(path)
    text, _confidence = _ocr_image(image)
    return text


def extract_text_from_file(path: str) -> str:
    extension = os.path.splitext(path)[1].lower()
    if extension == ".txt":
        return extract_text_from_txt(path)
    if extension == ".pdf":
        return extract_text_from_pdf(path)
    if extension == ".docx":
        return extract_text_from_docx(path)
    if extension == ".pptx":
        return extract_text_from_pptx(path)
    if extension in [".png", ".jpg", ".jpeg", ".bmp"]:
        return extract_text_from_image(path)
    raise ValueError("File type is not supported for extraction")


def extract_file(path: str, original_name: str | None = None) -> ExtractionResult:
    extension = os.path.splitext(path)[1].lower()
    file_name = original_name or os.path.basename(path)
    page_count = None
    ocr_pages = 0
    average_confidence = None

    if extension == ".txt":
        text = extract_text_from_txt(path)
        method = "plain_text"
    elif extension == ".pdf":
        text, page_count, ocr_pages, average_confidence, method = extract_pdf_with_ocr(path)
    elif extension == ".docx":
        text = extract_text_from_docx(path)
        method = "docx_text"
    elif extension == ".pptx":
        text = extract_text_from_pptx(path)
        method = "pptx_text"
    elif extension in [".png", ".jpg", ".jpeg", ".bmp", ".webp", ".tif", ".tiff"]:
        image = Image.open(path)
        text, average_confidence = _ocr_image(image)
        page_count = 1
        ocr_pages = 1 if text else 0
        method = "image_ocr"
    else:
        raise ValueError("File type is not supported for extraction")

    return ExtractionResult(
        text=text.strip(),
        file_name=file_name,
        file_type=extension.lstrip("."),
        extraction_method=method,
        page_count=page_count,
        ocr_pages=ocr_pages,
        average_confidence=average_confidence,
    )


def extract_text_from_upload(upload_file: UploadFile, target_dir: str) -> ExtractionResult:
    original_name = upload_file.filename
    saved_path = save_upload_file(upload_file, target_dir)
    try:
        extracted = extract_file(saved_path, original_name)
    except pytesseract.TesseractNotFoundError as exc:
        raise ValueError(
            "Tesseract OCR is not available. Install Tesseract or set TESSERACT_CMD in backend/.env."
        ) from exc
    except pytesseract.TesseractError as exc:
        raise ValueError(f"OCR failed: {exc}") from exc
    if not extracted.text:
        raise ValueError("No readable text could be extracted from this file")
    return extracted
